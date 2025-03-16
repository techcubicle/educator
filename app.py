from fastapi import FastAPI, UploadFile, File
import fitz  # PyMuPDF for PDF processing
import pptx  # python-pptx for PPT processing
import openai  # GPT-4 for explanations
import pytesseract  # OCR for image-to-text
from moviepy import VideoFileClip
import speech_recognition as sr  # Speech-to-text for MP4
import os

app = FastAPI()

OPENAI_API_KEY = "CjuFmMc06jdhIe2mLBIRdA8YxLDKJTnMbHEGwpdG3T3BlbkFJFMm8gnY3vVyXGYzWUSCPA7mB2s3upaycEN1XZBJ2A_EsvKIG3qIvZRKxLh9gr3hcbJiDsnDVMA"

# Function to extract text from PDFs
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = "\n".join([page.get_text() for page in doc])
    return text

# Function to extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = pptx.Presentation(ppt_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text

# Function to extract speech from MP4
def extract_audio_from_mp4(mp4_path):
    clip = VideoFileClip(mp4_path)
    audio_path = "audio.wav"
    clip.audio.write_audiofile(audio_path)
    return audio_path

# Function to transcribe speech from audio file
def transcribe_audio(audio_path):
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        audio = recognizer.record(source)
    try:
        return recognizer.recognize_google(audio)
    except sr.UnknownValueError:
        return "Could not understand audio"
    except sr.RequestError:
        return "Speech recognition service unavailable"

# Function to get AI-generated explanations
def get_ai_explanation(content):
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "system", "content": "Explain the following concept with real-world examples:"},
                  {"role": "user", "content": content}]
    )
    return response["choices"][0]["message"]["content"]

@app.post("/upload/pdf")
async def upload_pdf(file: UploadFile = File(...)):
    file_path = f"temp/{file.filename}"
    with open(file_path, "wb") as f:
        f.write(await file.read())
    
    text = extract_text_from_pdf(file_path)
    explanation = get_ai_explanation(text)
    
    os.remove(file_path)
    return {"original_text": text, "ai_explanation": explanation}

@app.post("/upload/ppt")
async def upload_ppt(file: UploadFile = File(...)):
    file_path = f"temp/{file.filename}"
    with open(file_path, "wb") as f:
        f.write(await file.read())
    
    text = extract_text_from_ppt(file_path)
    explanation = get_ai_explanation(text)
    
    os.remove(file_path)
    return {"original_text": text, "ai_explanation": explanation}

@app.post("/upload/mp4")
async def upload_mp4(file: UploadFile = File(...)):
    file_path = f"temp/{file.filename}"
    with open(file_path, "wb") as f:
        f.write(await file.read())
    
    audio_path = extract_audio_from_mp4(file_path)
    transcript = transcribe_audio(audio_path)
    explanation = get_ai_explanation(transcript)
    
    os.remove(file_path)
    os.remove(audio_path)
    
    return {"original_text": transcript, "ai_explanation": explanation}
